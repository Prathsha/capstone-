"""
context_builder.py
------------------
Assembles a rich, structured context string from all available data sources
and injects it into the Gemini system prompt.

HOW TO ADD A NEW DATA SOURCE
-----------------------------
1. Write an async function `_fetch_<source>(account, client)` that returns
   a dict or list of plain dicts.
2. Add a call to it inside `build_account_contexts()`.
3. Format its output into a readable section inside `_format_account_block()`.

The rest of the chat endpoint requires no changes.

STATIC FILE LOADERS
--------------------
_load_ibm_entitlements(), _load_competitive_installs(), _load_pipeline_opps(),
and _load_pptx_summary() read the EPM CSV files and the Data Platform PPTX from
backend/data/ and append their output to the context string unconditionally.
All loaders degrade gracefully — any missing file or library returns "".
"""

import asyncio
import os
from datetime import datetime, timedelta
from pathlib import Path

import httpx

from config import settings
from data import ACCOUNTS, SELLER
from queries import news_query

# ── Path constants ────────────────────────────────────────────────────────────

_DATA_DIR = Path(__file__).resolve().parent / "data"
_ACCOUNT_DIR = _DATA_DIR / "account"
_DOCS_DIR = _DATA_DIR / "docs"

_IBM_CSV      = _ACCOUNT_DIR / "US - HRZ - SEI INVESTMENTS_EPM_IBM_20260714.csv"
_COMP_CSV     = _ACCOUNT_DIR / "US - HRZ - SEI INVESTMENTS_EPM_Competitive_20260714.csv"
_OPP_CQ_CSV   = _ACCOUNT_DIR / "US - HRZ - SEI INVESTMENTS_EPM_Opportunity_CQ_20260714.csv"
_OPP_NQ_CSV   = _ACCOUNT_DIR / "US - HRZ - SEI INVESTMENTS_EPM_Opportunity_NQ_20260714.csv"
_PPTX_FILE    = _DOCS_DIR / "Data Platform 101.pptx"


# ── Per-account data fetchers ────────────────────────────────────────────────

async def _fetch_news(account: dict, client: httpx.AsyncClient) -> list[dict]:
    """Fetch recent NewsAPI articles for one account (max 5 headlines)."""
    if not settings.NEWS_API_KEY:
        return []
    query = news_query(account["name"])
    from_date = (datetime.utcnow() - timedelta(days=30)).strftime("%Y-%m-%d")
    try:
        resp = await client.get(
            "https://newsapi.org/v2/everything",
            params={
                "q": query,
                "from": from_date,
                "sortBy": "publishedAt",
                "language": "en",
                "pageSize": 5,
                "apiKey": settings.NEWS_API_KEY,
            },
        )
        if resp.status_code != 200:
            return []
        return [
            {"title": a["title"], "source": a["source"]["name"], "published": a["publishedAt"][:10]}
            for a in resp.json().get("articles", [])
            if a.get("title") and "[Removed]" not in a.get("title", "")
        ][:5]
    except Exception:
        return []


async def _fetch_stock(account: dict, client: httpx.AsyncClient) -> dict | None:
    """Fetch a Finnhub stock quote for publicly-traded accounts."""
    ticker = account.get("ticker")
    if not ticker or not settings.FINNHUB_API_KEY:
        return None
    try:
        resp = await client.get(
            "https://finnhub.io/api/v1/quote",
            params={"symbol": ticker, "token": settings.FINNHUB_API_KEY},
        )
        if resp.status_code != 200:
            return None
        d = resp.json()
        return {
            "ticker": ticker,
            "current_price": d.get("c"),
            "change_pct": round((d.get("dp") or 0), 2),
            "52w_high": d.get("h"),
            "52w_low": d.get("l"),
        }
    except Exception:
        return None


# ── Static file loaders ──────────────────────────────────────────────────────

def _load_ibm_entitlements() -> str:
    """Return a formatted section of IBM products installed at the account (from EPM CSV)."""
    try:
        import pandas as pd
        if not _IBM_CSV.exists():
            return ""
        df = pd.read_csv(_IBM_CSV, encoding="utf-16", sep="\t")
        col = "UT Lvl 30 Name"
        if col not in df.columns:
            return ""
        products = sorted(df[col].dropna().unique().tolist())
        if not products:
            return ""
        lines = ["## IBM Installed Base (from EPM)"]
        for p in products:
            lines.append(f"- {p}")
        return "\n".join(lines)
    except Exception:
        return ""


def _load_competitive_installs() -> str:
    """Return a formatted section of competitive products installed at the account."""
    try:
        import pandas as pd
        if not _COMP_CSV.exists():
            return ""
        df = pd.read_csv(_COMP_CSV, encoding="utf-16", sep="\t")
        if "Competitor Name" not in df.columns or "Competitor Product Name" not in df.columns:
            return ""
        # Deduplicate by vendor + product
        pairs = (
            df[["Competitor Name", "Competitor Product Name"]]
            .dropna()
            .drop_duplicates()
            .sort_values(["Competitor Name", "Competitor Product Name"])
        )
        if pairs.empty:
            return ""
        # Cap at 50 most relevant entries to stay within token budget
        pairs = pairs.head(50)
        lines = ["## Competitive Installs"]
        for _, row in pairs.iterrows():
            lines.append(f"- {row['Competitor Name']}: {row['Competitor Product Name']}")
        return "\n".join(lines)
    except Exception:
        return ""


def _load_pipeline_opps() -> str:
    """Return formatted pipeline opportunities from both CQ and NQ CSV files."""
    key_cols = [
        "Opp Name", "Oppty Value", "ISC Sales Stage Name",
        "Sales Forecast Category Shortname", "Reporting Product Family",
        "UT Lvl 30 Name Dynamic", "Opp Next Step (P2P)",
    ]

    def _parse_opps(path: Path, label: str) -> str:
        try:
            import pandas as pd
            if not path.exists():
                return ""
            df = pd.read_csv(path, encoding="utf-16", sep="\t")
            existing = [c for c in key_cols if c in df.columns]
            if "Opp Name" not in existing:
                return ""
            df = df[existing].dropna(subset=["Opp Name"]).drop_duplicates(subset=["Opp Name"])
            lines = [f"## Pipeline Opportunities ({label})"]
            for _, row in df.iterrows():
                value = f"${int(row['Oppty Value']):,}" if "Oppty Value" in row and not pd.isna(row["Oppty Value"]) else "N/A"
                stage = row.get("ISC Sales Stage Name", "")
                forecast = row.get("Sales Forecast Category Shortname", "")
                product = row.get("UT Lvl 30 Name Dynamic", row.get("Reporting Product Family", ""))
                notes = row.get("Opp Next Step (P2P)", "")
                notes_str = f" | Notes: {str(notes).strip()}" if notes and str(notes) != "nan" else ""
                lines.append(
                    f"- {row['Opp Name']} | {value} | Stage: {stage} | "
                    f"Forecast: {forecast} | Product: {product}{notes_str}"
                )
            return "\n".join(lines)
        except Exception:
            return ""

    cq = _parse_opps(_OPP_CQ_CSV, "CQ")
    nq = _parse_opps(_OPP_NQ_CSV, "NQ")
    parts = [p for p in [cq, nq] if p]
    return "\n\n".join(parts)


def _load_pptx_summary() -> str:
    """Return a capped text summary of the Data Platform 101 PPTX."""
    try:
        from pptx import Presentation
    except ImportError:
        return ""
    try:
        if not _PPTX_FILE.exists():
            return ""
        prs = Presentation(str(_PPTX_FILE))
        texts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
        full_text = "\n".join(texts)
        # Cap at 2000 characters to stay within token budget
        if len(full_text) > 2000:
            full_text = full_text[:2000] + "\n... [truncated]"
        return f"## IBM Data Platform Product Reference\n{full_text}"
    except Exception:
        return ""


# ── Format helpers ───────────────────────────────────────────────────────────

def _format_account_block(account: dict, news: list[dict], stock: dict | None) -> str:
    """Render one account's data as a readable text block for the LLM."""
    lines: list[str] = []

    # Basic profile
    lines.append(f"### Account: {account['name']}")
    lines.append(f"- Industry: {account['industry']}")
    lines.append(f"- Region: {account['region']}")
    lines.append(f"- Tier: {account['tier']}")
    lines.append(f"- Health Score: {account['health_score']}/100")
    lines.append(f"- Last Contact: {account['last_contact_days_ago']} days ago")
    lines.append(f"- Renewal Date: {account['renewal_date']}")

    # Quota/pipeline
    quota = account["quota"]
    closed = account["closed"]
    pipeline = account["pipeline"]
    attainment = round(closed / quota * 100, 1) if quota else 0
    gap = quota - closed
    lines.append(f"- Quota: ${quota:,}  |  Closed Won: ${closed:,} ({attainment}%)  |  Pipeline: ${pipeline:,}  |  Gap to Quota: ${gap:,}")

    # IBM products
    lines.append(f"- Owned IBM Products: {', '.join(account.get('owned_products', []))}")

    # Pending actions
    pending = account.get("pending_actions", [])
    if pending:
        lines.append("- Pending Actions:")
        for p in pending:
            lines.append(f"  • [{p['priority'].upper()}] {p['type']}: {p['description']} (due {p['due_date']})")

    # Existing AI suggestions
    suggested = account.get("suggested_actions", [])
    if suggested:
        lines.append("- Existing AI Suggestions:")
        for s in suggested:
            lines.append(f"  • [{s['type']}] {s['description']} (confidence: {round(s['confidence']*100)}%)")

    # Stock data
    if stock:
        lines.append(
            f"- Stock ({stock['ticker']}): ${stock['current_price']} "
            f"({'+' if stock['change_pct'] >= 0 else ''}{stock['change_pct']}% today) | "
            f"52W range: ${stock['52w_low']}–${stock['52w_high']}"
        )
    elif account.get("ticker"):
        lines.append(f"- Stock ticker: {account['ticker']} (live data unavailable)")
    else:
        lines.append("- Privately held (no public stock data)")

    # Recent news
    if news:
        lines.append("- Recent News Headlines:")
        for n in news:
            lines.append(f"  • [{n['published']}] {n['source']}: {n['title']}")
    else:
        lines.append("- Recent News: (no live news data — NEWS_API_KEY not configured)")

    return "\n".join(lines)


# ── Public API ───────────────────────────────────────────────────────────────

async def build_context(account_ids: list[str] | None = None) -> str:
    """
    Fetch all available data for the requested accounts concurrently and
    return a fully formatted context string ready to embed in the system prompt.

    Pass account_ids=None to include all accounts.
    """
    if account_ids:
        accounts = [a for a in ACCOUNTS if a["id"] in account_ids]
    else:
        accounts = ACCOUNTS

    # Concurrently fetch news + stock for every account
    async with httpx.AsyncClient(timeout=12.0) as client:
        tasks = [
            asyncio.gather(
                _fetch_news(account, client),
                _fetch_stock(account, client),
            )
            for account in accounts
        ]
        results = await asyncio.gather(*tasks)

    # Build per-account blocks
    account_blocks = [
        _format_account_block(account, news, stock)
        for account, (news, stock) in zip(accounts, results)
    ]

    # Seller summary
    seller_lines = [
        "## Seller Summary",
        f"- Name: {SELLER['name']} | Title: {SELLER['title']} | Pod: {SELLER['pod']}",
        f"- Total Quota: ${SELLER['quota_total']:,}",
        f"- Closed Won: ${SELLER['quota_closed']:,} ({round(SELLER['quota_closed']/SELLER['quota_total']*100,1)}%)",
        f"- Pipeline: ${SELLER['quota_pipeline']:,}",
        f"- YTD Target: {round(SELLER['ytd_target_pct']*100)}%",
    ]

    # Static file sections (EPM CSVs + PPTX) — account-agnostic, always appended
    static_sections = [
        _load_ibm_entitlements(),
        _load_competitive_installs(),
        _load_pipeline_opps(),
        _load_pptx_summary(),
    ]
    static_parts = [s for s in static_sections if s]

    all_parts = ["\n".join(seller_lines)] + account_blocks + static_parts
    return "\n\n".join(all_parts)
